"""Build .pptx fixtures: modify-protected and open.

python-pptx cannot write a modify password, so the <p:modifyVerifier> element
is injected into ppt/presentation.xml directly -- which is exactly what
PowerPoint writes when you set "password to modify".

    python tests/make_pptx_fixture.py out.pptx [--open]
"""

import re
import shutil
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

MODIFY_VERIFIER = (
    '<p:modifyVerifier cryptProviderType="rsaAES" cryptAlgorithmClass="hash" '
    'cryptAlgorithmType="typeAny" cryptAlgorithmSid="14" spinCount="100000" '
    'saltData="abcdefghijklmnopqrstuvw==" hashData="Zm9vYmFyYmF6cXV4"/>'
)


def _base(path: Path, title: str) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    box.text_frame.text = "Second line of the slide body."

    second = prs.slides.add_slide(prs.slide_layouts[5])
    second.shapes.title.text = "Slide two"

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return path


def build_open(path: Path) -> Path:
    """No protection."""
    return _base(path, "This presentation has no protection.")


def build_protected(path: Path) -> Path:
    """Carries a modify password, so PowerPoint opens it read-only."""
    _base(path, "This presentation is modify-protected.")

    # Inject the element the way PowerPoint does, preserving entry order.
    tmp = path.with_suffix(".tmp")
    with zipfile.ZipFile(path, "r") as src, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as out:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename.lower() == "ppt/presentation.xml":
                xml = data.decode("utf-8")
                # modifyVerifier belongs just before </p:presentation>.
                xml = re.sub(r"</p:presentation>", MODIFY_VERIFIER + "</p:presentation>", xml, count=1)
                data = xml.encode("utf-8")
            info = zipfile.ZipInfo(item.filename, date_time=item.date_time)
            info.compress_type = item.compress_type
            info.external_attr = item.external_attr
            out.writestr(info, data)

    shutil.move(str(tmp), str(path))
    return path


if __name__ == "__main__":
    out = Path(sys.argv[1] if len(sys.argv) > 1 else "tests/fixtures/protected.pptx")
    if "--open" in sys.argv:
        build_open(out)
    else:
        build_protected(out)
    print(f"wrote {out} ({out.stat().st_size} bytes)")
