"""Tests for .pptx modify-protection removal and macro-enabled variants."""

import shutil
import sys
import zipfile
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from backend import dispatch  # noqa: E402
from backend.pptx_cracker import PptxCracker, inspect_pptx  # noqa: E402
from tests.make_fixture import build as build_docx  # noqa: E402
from tests.make_pptx_fixture import build_open, build_protected  # noqa: E402
from tests.make_xlsx_fixture import build_protected as build_xlsx  # noqa: E402


@pytest.fixture
def protected(tmp_path):
    return build_protected(tmp_path / "protected.pptx")


@pytest.fixture
def unprotected(tmp_path):
    return build_open(tmp_path / "open.pptx")


def has_verifier(path) -> bool:
    with zipfile.ZipFile(path) as zf:
        return b"modifyVerifier" in zf.read("ppt/presentation.xml")


# ---------------------------------------------------------------- core


def test_removes_modify_protection(protected, tmp_path):
    assert has_verifier(protected)

    result = PptxCracker(str(protected), str(tmp_path / "out")).unlock()
    assert result.status == "success", result.error
    assert result.file_format == "pptx"
    assert result.protections_found >= 1
    assert not has_verifier(result.output_path)


def test_original_is_untouched(protected, tmp_path):
    before = protected.read_bytes()
    PptxCracker(str(protected), str(tmp_path / "out")).unlock()
    assert protected.read_bytes() == before


def test_slide_content_survives(protected, tmp_path):
    from pptx import Presentation

    result = PptxCracker(str(protected), str(tmp_path / "out")).unlock()
    prs = Presentation(result.output_path)
    assert len(prs.slides) == 2
    titles = [s.shapes.title.text for s in prs.slides if s.shapes.title]
    assert "This presentation is modify-protected." in titles
    assert "Slide two" in titles


def test_entry_order_is_preserved(protected, tmp_path):
    result = PptxCracker(str(protected), str(tmp_path / "out")).unlock()
    with zipfile.ZipFile(protected) as src, zipfile.ZipFile(result.output_path) as out:
        assert src.namelist() == out.namelist()
        assert out.testzip() is None


def test_unrelated_parts_are_byte_identical(protected, tmp_path):
    result = PptxCracker(str(protected), str(tmp_path / "out")).unlock()
    with zipfile.ZipFile(protected) as src, zipfile.ZipFile(result.output_path) as out:
        for name in src.namelist():
            if name.lower() == "ppt/presentation.xml":
                continue
            assert src.read(name) == out.read(name), f"{name} changed"


def test_unprotected_presentation_still_succeeds(unprotected, tmp_path):
    result = PptxCracker(str(unprotected), str(tmp_path / "out")).unlock()
    assert result.status == "success"
    assert not has_verifier(result.output_path)


def test_rejects_password_protected_presentation(tmp_path):
    p = tmp_path / "encrypted.pptx"
    p.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 512)
    result = PptxCracker(str(p)).unlock()
    assert result.status == "failed"
    assert "password-protected" in result.error


def test_rejects_zip_without_a_presentation(tmp_path):
    p = tmp_path / "notppt.pptx"
    with zipfile.ZipFile(p, "w") as zf:
        zf.writestr("hello.txt", "nope")
    result = PptxCracker(str(p)).unlock()
    assert result.status == "failed"
    assert "presentation.xml" in result.error


def test_logs_cover_all_seven_steps(protected, tmp_path):
    result = PptxCracker(str(protected), str(tmp_path / "out")).unlock()
    joined = "\n".join(result.logs)
    for n in range(1, 8):
        assert f"[{n}/7]" in joined


def test_no_temp_directory_leaks(protected, tmp_path):
    cracker = PptxCracker(str(protected), str(tmp_path / "out"))
    cracker.unlock()
    assert not cracker._temp_dir.exists()


def test_inspect_reports_modify_protection(protected, unprotected):
    assert inspect_pptx(str(protected))["protected"] is True
    assert inspect_pptx(str(unprotected))["protected"] is False


# ------------------------------------------------- macro-enabled variants


def test_pptm_is_handled_like_pptx(protected, tmp_path):
    """A .pptm is a .pptx with macros -- same parts, same protection."""
    macro = tmp_path / "deck.pptm"
    shutil.copy2(protected, macro)

    result = dispatch.unlock(str(macro), str(tmp_path / "out"))
    assert result.status == "success", result.error
    assert result.file_format == "pptm"
    # The extension must survive, or the file stops being macro-enabled.
    assert Path(result.output_path).suffix == ".pptm"
    assert not has_verifier(result.output_path)


def test_docm_is_handled_like_docx(tmp_path):
    src = build_docx(tmp_path / "d.docx", protected=True)
    macro = tmp_path / "notes.docm"
    shutil.copy2(src, macro)

    result = dispatch.unlock(str(macro), str(tmp_path / "out"))
    assert result.status == "success", result.error
    assert result.file_format == "docm"
    assert Path(result.output_path).suffix == ".docm"
    with zipfile.ZipFile(result.output_path) as zf:
        assert b"documentProtection" not in zf.read("word/settings.xml")


def test_xlsm_is_handled_like_xlsx(tmp_path):
    src = build_xlsx(tmp_path / "b.xlsx")
    macro = tmp_path / "book.xlsm"
    shutil.copy2(src, macro)

    result = dispatch.unlock(str(macro), str(tmp_path / "out"))
    assert result.status == "success", result.error
    assert result.file_format == "xlsm"
    assert Path(result.output_path).suffix == ".xlsm"


# ------------------------------------------------------------ dispatch


def test_dispatch_knows_every_office_format():
    for ext, label in [
        (".docx", "DOCX"), (".docm", "DOCM"),
        (".xlsx", "XLSX"), (".xlsm", "XLSM"),
        (".pptx", "PPTX"), (".pptm", "PPTM"),
        (".pdf", "PDF"),
    ]:
        assert dispatch.format_label(f"a{ext}") == label
        assert dispatch.detect_format(f"a{ext}") == ext.lstrip(".")


def test_every_format_has_a_step_count():
    for ext in dispatch.SUPPORTED:
        assert ext.lstrip(".") in dispatch.STEP_COUNT


def test_dispatch_routes_pptx(protected, tmp_path):
    result = dispatch.unlock(str(protected), str(tmp_path / "out"))
    assert result.status == "success"
    assert result.file_format == "pptx"
    assert result.method == "ooxml"


def test_dispatch_inspect_pptx(protected):
    assert dispatch.inspect(str(protected))["format"] == "pptx"
